"""
Deep Learning Presentation Generator
生成完整的深度学习分享PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ===== Slide 1: 标题页 =====
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "深度学习入门：从理论到实践"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Neural Networks · Training · Forward/Backward Propagation"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.alignment = PP_ALIGN.CENTER
    
    # ===== Slide 2: 目录 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "目录")
    
    content = """1. 神经网络基础
   • 生物神经元启发
   • 神经网络结构（输入层/隐藏层/输出层）

2. 学习与训练核心概念
   • 线性回归 → 逻辑回归
   • 损失函数、权重与偏移
   • 梯度下降与链式求导
   • 欠拟合与过拟合

3. 数据集与学习方法
   • 训练集/测试集/交叉验证集
   • 监督学习 vs 非监督学习

4. 正向传播与反向传播

5. 实践演示
   • 实验一：简单的深度学习入门
   • 实验二：人脸识别实战"""
    
    add_content(slide, content, Inches(1.5), Inches(1.8))
    
    # ===== Slide 3: 神经网络简介 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "神经网络简介：从生物到计算")
    
    content = """生物神经元 (Biological Neuron)
┌─────────────────────────────────────────┐
│  树突(Dendrites) → 胞体(Soma) → 轴突(Axon) → 突触(Synapse) │
│  接收信号      →  处理    →  输出    →  传递给下一个神经元   │
└─────────────────────────────────────────┘

人工神经元 (Artificial Neuron) - McCulloch-Pitts 模型 (1943)
┌─────────────────────────────────────────┐
│  输入 x₁, x₂, ..., xₙ                    │
│    ↓                                    │
│  加权求和: z = w₁x₁ + w₂x₂ + ... + wₙxₙ + b │
│    ↓                                    │
│  激活函数: a = g(z)                      │
│    ↓                                    │
│  输出                                    │
└─────────────────────────────────────────┘

• 权重 (Weights, w): 表示输入的重要性
• 偏移/偏置 (Bias, b): 调整激活阈值"""
    
    add_content(slide, content, Inches(0.8), Inches(1.5))
    
    # ===== Slide 4: 神经网络结构 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "神经网络结构：三层架构")
    
    content = """                    输入层          隐藏层          输出层
                    (Input)       (Hidden)       (Output)
                      
    x₁ ───────→    [○]            [○]            [○]  → ŷ (预测值)
                    ↓            / ↓ \            ↑
    x₂ ───────→    [○]    →    [○]──[○]    →    [○]
                    ↓            \ ↑ /            ↑
    x₃ ───────→    [○]            [○]            [○]
                    
    
• 输入层: 接收原始数据特征 (x₁, x₂, x₃...)
• 隐藏层: 提取特征、学习表示（可有多层 → "深度"学习）
• 输出层: 产生最终预测结果

深度 = 隐藏层数量（不含输入输出层）"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 5: 线性回归 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "从线性回归开始：预测房价")
    
    content = """场景: 根据房屋面积预测价格

数据样例:
┌──────────────┬──────────────┐
│ 面积(m²) x   │ 价格(万) y   │
├──────────────┼──────────────┤
│     50       │     100      │
│     80       │     160      │
│    100       │     200      │
│    120       │     240      │
└──────────────┴──────────────┘

线性模型:  ŷ = wx + b

• ŷ (y-hat): 预测值
• y: 真实目标值
• w: 权重 (Weight) - 每平方米价格
• b: 偏移 (Bias) - 基础价格

目标: 找到最优的 w 和 b，使得预测值 ŷ 尽可能接近真实值 y"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 6: 损失函数 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "损失函数：衡量预测的好坏")
    
    content = """为什么需要损失函数？
量化预测值 ŷ 与真实值 y 之间的差距

均方误差 (Mean Squared Error, MSE)
┌───────────────────────────────────────────────┐
│                                               │
│      1     m                                  │
│  J = ───   Σ  (ŷ⁽ⁱ⁾ - y⁽ⁱ⁾)²                  │
│      2m   i=1                                 │
│                                               │
└───────────────────────────────────────────────┘

• m: 样本数量
• ŷ⁽ⁱ⁾: 第 i 个样本的预测值
• y⁽ⁱ⁾: 第 i 个样本的真实值
• J: 代价函数 (Cost Function)

计算示例:
样本1: x=50, y=100, w=2, b=10
  ŷ = 2×50 + 10 = 110
  误差 = (110 - 100)² = 100

目标: 最小化 J(w, b)"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 7: 梯度下降 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "梯度下降：寻找最优参数")
    
    content = """核心思想: 沿着损失函数下降最快的方向更新参数

梯度下降算法:
┌───────────────────────────────────────────────┐
│                                               │
│  重复直到收敛:                                 │
│    w := w - α · ∂J/∂w   (对 w 求偏导)        │
│    b := b - α · ∂J/∂b   (对 b 求偏导)        │
│                                               │
└───────────────────────────────────────────────┘

• α (alpha): 学习率 (Learning Rate)
  - 太大: 可能震荡无法收敛
  - 太小: 收敛速度太慢
  - 典型值: 0.001, 0.01, 0.1

• ∂J/∂w: 损失函数对 w 的梯度（斜率）
  表示 w 变化时 J 的变化率

直观理解:
想象在山坡上往下走，梯度告诉你哪个方向最陡
每次迈一小步（学习率控制步长），最终到达山谷（最小值）"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 8: 链式求导 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "链式求导：反向传播的核心")
    
    content = """对于复合函数，梯度如何计算？

线性回归的梯度推导:
J = (1/2m) Σ(ŷ - y)²,  其中 ŷ = wx + b

∂J/∂w = ∂J/∂ŷ · ∂ŷ/∂w
      = (1/m) Σ(ŷ - y) · x

∂J/∂b = ∂J/∂ŷ · ∂ŷ/∂b  
      = (1/m) Σ(ŷ - y)

示例计算 (单样本):
x = 50, y = 100, w = 2, b = 10
ŷ = 2×50 + 10 = 110

∂J/∂w = (110 - 100) × 50 = 500
∂J/∂b = (110 - 100) = 10

若 α = 0.001:
w_new = 2 - 0.001×500 = 1.5
b_new = 10 - 0.001×10 = 9.99"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 9: 逻辑回归 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "逻辑回归：从回归到分类")
    
    content = """场景: 判断邮件是否为垃圾邮件（二分类问题）

Sigmoid 激活函数:
┌───────────────────────────────────────────────┐
│                                               │
│              1                                │
│  σ(z) = ──────────                            │
│         1 + e⁻ᶻ                               │
│                                               │
└───────────────────────────────────────────────┘

性质:
• 输出范围 (0, 1) → 可解释为概率
• z = 0 时, σ(z) = 0.5
• z → +∞, σ(z) → 1
• z → -∞, σ(z) → 0

逻辑回归模型:
  z = wx + b
  ŷ = σ(z) = P(y=1|x)

损失函数 (交叉熵):
  L = -[y·log(ŷ) + (1-y)·log(1-ŷ)]"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 10: 拟合问题 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "欠拟合与过拟合")
    
    content = """                    欠拟合                          过拟合
                 (Underfitting)                  (Overfitting)
                    
    数据分布:         ●  ●    ●                      ●  ●    ●
                   ●        ●  ●                  ●        ●  ●
                 ●    ●        ●                ●    ●        ●
                                                  ●          ●
                                                    ●      ●
                    
    模型结果:    简单直线拟合                  复杂曲线穿过每个点
                训练误差高                     训练误差很低
                测试误差也高                   测试误差很高
                
    原因:        模型太简单                    模型太复杂/数据太少
    
    解决方案:
    • 欠拟合: 增加特征、使用更复杂模型、减少正则化
    • 过拟合: 增加数据、正则化(L1/L2)、Dropout、早停"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 11: 数据集划分 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "数据集划分：训练/验证/测试")
    
    content = """为什么需要划分？
评估模型泛化能力（对未见过数据的预测能力）

┌─────────────────────────────────────────────────────────────┐
│                                                             │
│   原始数据集                                                │
│   ┌─────────────────────────────────────────────────────┐   │
│   │                                                     │   │
│   │  ┌───────────────┐  ┌─────────────┐  ┌──────────┐  │   │
│   │  │   训练集      │  │  验证集     │  │  测试集  │  │   │
│   │  │   (70%)       │  │  (15%)      │  │  (15%)   │  │   │
│   │  │               │  │             │  │          │  │   │
│   │  │ • 训练模型     │  │ • 调参      │  │ • 最终   │  │   │
│   │  │ • 学习权重     │  │ • 选模型    │  │   评估   │  │   │
│   │  │ • 梯度下降     │  │ • 早停      │  │ • 不用于 │  │   │
│   │  │               │  │             │  │   调参   │  │   │
│   │  └───────────────┘  └─────────────┘  └──────────┘  │   │
│   │                                                     │   │
│   └─────────────────────────────────────────────────────┘   │
│                                                             │
└─────────────────────────────────────────────────────────────┘

交叉验证 (K-Fold): 将数据分成K份，轮流用K-1份训练，1份验证"""
    
    add_content(slide, content, Inches(0.5), Inches(1.3))
    
    # ===== Slide 12: 监督 vs 非监督学习 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "监督学习 vs 非监督学习")
    
    content = """┌─────────────────────┬─────────────────────┐
│     监督学习        │     非监督学习       │
├─────────────────────┼─────────────────────┤
│                     │                     │
│  数据: (X, y)       │  数据: (X)          │
│  有标签             │  无标签             │
│                     │                     │
│  目标: 学习 X→y     │  目标: 发现数据     │
│        的映射       │        的内在结构   │
│                     │                     │
├─────────────────────┼─────────────────────┤
│  典型任务:          │  典型任务:          │
│  • 分类             │  • 聚类             │
│  • 回归             │  • 降维             │
│  • 目标检测         │  • 异常检测         │
│                     │                     │
├─────────────────────┼─────────────────────┤
│  示例:              │  示例:              │
│  • 垃圾邮件分类     │  • 客户分群         │
│  • 房价预测         │  • 主题建模         │
│  • 图像识别         │  • 数据压缩         │
│                     │                     │
└─────────────────────┴─────────────────────┘"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 13: 正向传播 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "正向传播：计算预测值")
    
    content = """正向传播 (Forward Propagation): 从输入到输出的计算流程

网络结构:
  输入层        隐藏层1         隐藏层2        输出层
    [x₁]         [h₁⁽¹⁾]        [h₁⁽²⁾]         [ŷ]
    [x₂]    →    [h₂⁽¹⁾]   →    [h₂⁽²⁾]   →     
    [x₃]         [h₃⁽¹⁾]        [h₃⁽²⁾]

计算过程:
Layer 1:  z⁽¹⁾ = W⁽¹⁾x + b⁽¹⁾,   a⁽¹⁾ = g(z⁽¹⁾)
Layer 2:  z⁽²⁾ = W⁽²⁾a⁽¹⁾ + b⁽²⁾, a⁽²⁾ = g(z⁽²⁾)
Output:   z⁽³⁾ = W⁽³⁾a⁽²⁾ + b⁽³⁾, ŷ = σ(z⁽³⁾)

维度检查:
• x: (n_x, 1)
• W⁽¹⁾: (n_h1, n_x)
• W⁽²⁾: (n_h2, n_h1)
• W⁽³⁾: (n_y, n_h2)

一次正向传播 = 一次预测"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 14: 反向传播 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "反向传播：计算梯度并更新参数")
    
    content = """反向传播 (Backpropagation): 从输出到输入的梯度传递

核心: 链式求导法则的应用

计算流程 (从右到左):

1. 输出层梯度:
   δ⁽ᴸ⁾ = ∂J/∂z⁽ᴸ⁾ = ŷ - y

2. 反向传播到隐藏层:
   δ⁽ˡ⁾ = (W⁽ˡ⁺¹⁾ᵀ · δ⁽ˡ⁺¹⁾) ⊙ g'(z⁽ˡ⁾)
   
3. 计算参数梯度:
   ∂J/∂W⁽ˡ⁾ = δ⁽ˡ⁾ · a⁽ˡ⁻¹⁾ᵀ
   ∂J/∂b⁽ˡ⁾ = δ⁽ˡ⁾

4. 更新参数:
   W⁽ˡ⁾ := W⁽ˡ⁾ - α · ∂J/∂W⁽ˡ⁾
   b⁽ˡ⁾ := b⁽ˡ⁾ - α · ∂J/∂b⁽ˡ⁾

关键: 梯度从输出层逐层反向传递到输入层，高效计算所有参数梯度"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 15: 实验一介绍 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验一：最简单的深度学习实践")
    
    content = """目标: 用最基本的API实现一个神经网络，理解核心流程

任务: 二分类问题（如判断点是红还是蓝）

实验内容:
1. 使用 NumPy 手动实现:
   • 网络初始化
   • 正向传播
   • 损失计算
   • 反向传播
   • 参数更新

2. 对比使用 PyTorch/Keras:
   • 同样的网络结构
   • 高层API封装
   • 理解框架在背后做了什么

代码路径: experiments/experiment_1_basic_nn/

预计时间: 15-20 分钟演示"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 16: 实验二介绍 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验二：人脸识别实战")
    
    content = """目标: 从训练到实时预测，完整体验深度学习工作流

场景: 识别照片中的人是谁

流程:
1. 数据准备
   • 收集少量照片（每人10-20张）
   • 数据增强（旋转、翻转、缩放）

2. 模型训练
   • 使用预训练模型（Transfer Learning）
   • 或在预训练模型上 Fine-tune

3. 实时识别
   • 现场拍照
   • 模型推理预测
   • 显示识别结果

代码路径: experiments/experiment_2_face_recognition/

预计时间: 15-20 分钟演示"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 17: 总结 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总结：深度学习的核心要素")
    
    content = """1. 网络结构
   • 层数、神经元数、连接方式
   • 激活函数的选择

2. 损失函数
   • 定义优化目标
   • 回归:MSE, 分类:Cross-Entropy

3. 优化算法
   • 梯度下降及其变种
   • 学习率调度

4. 数据
   • 质量与数量
   • 划分策略

5. 训练技巧
   • 正则化防止过拟合
   • Batch Normalization
   • 合适的初始化

"深度学习没有魔法，只有数学和工程。"""
    
    add_content(slide, content, Inches(0.8), Inches(1.3))
    
    # ===== Slide 18: Q&A =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Q & A")
    
    content = """感谢聆听！

代码和数据已准备好，接下来进行实践演示。


推荐学习资源:
• Deep Learning Specialization (Andrew Ng)
• 《深度学习》(花书)
• PyTorch/TensorFlow 官方教程
• Papers With Code"""
    
    add_content(slide, content, Inches(2), Inches(2.5))
    
    return prs

def add_title(slide, text):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)

def add_content(slide, text, left, top):
    """添加内容文本"""
    content_box = slide.shapes.add_textbox(left, top, Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.name = "Courier New" if "┌" in text or "│" in text else "微软雅黑"
    p.font.color.rgb = RGBColor(51, 51, 51)

if __name__ == "__main__":
    prs = create_presentation()
    prs.save("/Users/tuyouwu/sharing/deep_learning/slides/deep_learning_intro.pptx")
    print("PPT 生成成功: /Users/tuyouwu/sharing/deep_learning/slides/deep_learning_intro.pptx")
